"""Build Smart Dental Desk pitch deck as a real editable .pptx file.

Run:   python funding/build_deck.py
Output: funding/SmartDentalDesk-PitchDeck.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
import os

# ── Palette ─────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0B, 0x12, 0x20)   # deep navy
PANEL    = RGBColor(0x0F, 0x1A, 0x2E)   # card / panel
INK      = RGBColor(0xE7, 0xED, 0xF5)   # primary text
MUTED    = RGBColor(0x9F, 0xB0, 0xC8)   # secondary text
ACCENT   = RGBColor(0x4E, 0xA1, 0xFF)   # blue
ACCENT2  = RGBColor(0x7A, 0xFC, 0xC0)   # green (highlight numbers)
BORDER   = RGBColor(0x1B, 0x2A, 0x44)   # card border

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)  # 16:9

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # blank


# ── Helpers ─────────────────────────────────────────────────────────────────
def add_slide():
    s = prs.slides.add_slide(blank_layout)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.shadow.inherit = False
    return s


def add_textbox(slide, left, top, width, height, text, *,
                size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, font="Calibri", line_spacing=1.15):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor

    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_rich(slide, left, top, width, height, runs, *,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    """runs: list of (text, dict(size=, bold=, color=, font=))"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for text, opts in runs:
        if text == "\n":
            p = tf.add_paragraph()
            p.alignment = align
            p.line_spacing = line_spacing
            continue
        r = p.add_run()
        r.text = text
        r.font.name = opts.get("font", "Calibri")
        r.font.size = Pt(opts.get("size", 16))
        r.font.bold = opts.get("bold", False)
        r.font.color.rgb = opts.get("color", INK)
    return tb


def add_card(slide, left, top, width, height, *, fill=PANEL, border=BORDER):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.adjustments[0] = 0.06
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_pill(slide, left, top, text):
    w, h = Inches(3.2), Inches(0.35)
    p = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    p.adjustments[0] = 0.5
    p.fill.solid()
    p.fill.fore_color.rgb = RGBColor(0x13, 0x2A, 0x4A)
    p.line.color.rgb = ACCENT
    p.line.width = Pt(0.5)
    p.shadow.inherit = False
    tf = p.text_frame
    tf.margin_left = tf.margin_right = Inches(0.1)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.word_wrap = False
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    r = para.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = ACCENT
    return p


def add_chip(slide, left, top, text):
    w = Inches(max(1.4, 0.13 * len(text) + 0.4))
    h = Inches(0.34)
    p = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    p.adjustments[0] = 0.5
    p.fill.solid()
    p.fill.fore_color.rgb = RGBColor(0x14, 0x22, 0x3D)
    p.line.color.rgb = BORDER
    p.line.width = Pt(0.5)
    p.shadow.inherit = False
    tf = p.text_frame
    tf.margin_left = tf.margin_right = Inches(0.12)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.word_wrap = False
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    r = para.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(11)
    r.font.color.rgb = INK
    return w


def add_footer(slide, num, total=14):
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(8), Inches(0.3),
                "SMART DENTAL DESK  ·  Investor Pitch  ·  2026",
                size=10, color=MUTED, bold=True)
    add_textbox(slide, Inches(12.0), Inches(7.05), Inches(1.0), Inches(0.3),
                f"{num:02d} / {total:02d}",
                size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def section_label(slide, text, left=Inches(0.6), top=Inches(0.55)):
    add_textbox(slide, left, top, Inches(8), Inches(0.35),
                text.upper(), size=12, bold=True, color=ACCENT)


def slide_title(slide, text, left=Inches(0.6), top=Inches(0.95), width=Inches(12), size=34):
    add_textbox(slide, left, top, width, Inches(1.4),
                text, size=size, bold=True, color=INK, line_spacing=1.1)


# ── KPI tile (for stat cards) ───────────────────────────────────────────────
def kpi_tile(slide, left, top, width, height, label, value, sub, *, val_color=ACCENT2):
    add_card(slide, left, top, width, height)
    add_textbox(slide, left + Inches(0.25), top + Inches(0.2),
                width - Inches(0.5), Inches(0.3),
                label.upper(), size=10, bold=True, color=MUTED)
    add_textbox(slide, left + Inches(0.25), top + Inches(0.55),
                width - Inches(0.5), Inches(0.9),
                value, size=34, bold=True, color=val_color)
    add_textbox(slide, left + Inches(0.25), top + Inches(1.45),
                width - Inches(0.5), height - Inches(1.6),
                sub, size=11, color=MUTED, line_spacing=1.25)


# Bullet card: label + bullet points
def bullet_card(slide, left, top, width, height, label, bullets, *,
                label_color=ACCENT, bullet_color=INK, body_size=12):
    add_card(slide, left, top, width, height)
    add_textbox(slide, left + Inches(0.25), top + Inches(0.2),
                width - Inches(0.5), Inches(0.3),
                label.upper(), size=10, bold=True, color=label_color)
    tb = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.55),
                                  width - Inches(0.5), height - Inches(0.7))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.25
        p.space_after = Pt(3)
        r = p.add_run()
        r.text = "•  " + b
        r.font.name = "Calibri"
        r.font.size = Pt(body_size)
        r.font.color.rgb = bullet_color


def paragraph_card(slide, left, top, width, height, label, body, *,
                   label_color=ACCENT, body_color=INK, body_size=13):
    add_card(slide, left, top, width, height)
    add_textbox(slide, left + Inches(0.25), top + Inches(0.2),
                width - Inches(0.5), Inches(0.3),
                label.upper(), size=10, bold=True, color=label_color)
    add_textbox(slide, left + Inches(0.25), top + Inches(0.55),
                width - Inches(0.5), height - Inches(0.7),
                body, size=body_size, color=body_color, line_spacing=1.3)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Cover
# ──────────────────────────────────────────────────────────────────────────────
s = add_slide()
add_pill(s, Inches(0.6), Inches(0.55), "PRE-SEED / SEED   ·   HEALTHCARE AI")

add_rich(s, Inches(0.6), Inches(1.2), Inches(10), Inches(2.6), [
    ("The AI-Powered Operating", {"size": 50, "bold": True, "color": INK}),
    ("\n", {}),
    ("System for Dental Practices.", {"size": 50, "bold": True, "color": INK}),
], line_spacing=1.05)

add_textbox(s, Inches(0.6), Inches(4.05), Inches(11), Inches(1.2),
            "Cloud-native, AI-first practice management — built to cut admin workload by 40%, "
            "no-shows by 30%, and bring real AI to every dental chair.",
            size=18, color=MUTED, line_spacing=1.35)

# Chip row
chips = ["7 AI features live", "Web · Mobile · API", "Multi-tenant SaaS",
         "Pilot revenue in India", "Canada-ready architecture"]
x = Inches(0.6)
for c in chips:
    w = add_chip(s, x, Inches(5.6), c)
    x += w + Inches(0.12)

# Right-side ask block (kept clear of title)
add_textbox(s, Inches(9.6), Inches(1.2), Inches(3.2), Inches(0.3),
            "RAISING", size=11, bold=True, color=MUTED, align=PP_ALIGN.RIGHT)
add_textbox(s, Inches(9.6), Inches(1.55), Inches(3.2), Inches(1.0),
            "CAD $[TBD]", size=38, bold=True, color=ACCENT2, align=PP_ALIGN.RIGHT)
add_textbox(s, Inches(9.6), Inches(2.5), Inches(3.2), Inches(0.6),
            "Pre-seed / Seed\n[TBD] month runway",
            size=12, color=MUTED, align=PP_ALIGN.RIGHT, line_spacing=1.3)

add_footer(s, 1)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Problem
# ──────────────────────────────────────────────────────────────────────────────
s = add_slide()
section_label(s, "The Problem")
slide_title(s, "Dental practices run on software built\nbefore the cloud — and before AI.")

col_w = Inches(6.0); col_h = Inches(2.1)
gap = Inches(0.25)
left1 = Inches(0.6); left2 = left1 + col_w + gap
row1 = Inches(3.0);  row2 = row1 + col_h + gap

paragraph_card(s, left1, row1, col_w, col_h, "Admin overload",
               "Front-desk teams spend 30–45% of their day on appointments, recalls, "
               "insurance forms, and follow-up — repetitive, rules-based work.")
paragraph_card(s, left2, row1, col_w, col_h, "No-show economics",
               "A typical 3-operatory clinic loses CAD $[TBD]/year to no-shows alone. "
               "Most clinics have no system that flags high-risk appointments.")
paragraph_card(s, left1, row2, col_w, col_h, "Documentation burnout",
               "Dentists spend 30+ minutes/day retyping the same SOAP notes and "
               "prescriptions. Top-3 driver of mid-career burnout in CDA surveys.")
paragraph_card(s, left2, row2, col_w, col_h, "Stagnant incumbents",
               "Dentrix, ABELDent, ClearDent — pre-2010 architectures, on-prem, "
               "bolted-on AI, priced for chains, hostile to single-chair clinics.")
add_footer(s, 2)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Solution
# ──────────────────────────────────────────────────────────────────────────────
s = add_slide()
section_label(s, "Our Solution")
slide_title(s, "One integrated platform. Three surfaces.\nSeven live AI features.")

col_w = Inches(4.0); col_h = Inches(2.0); gap = Inches(0.2)
left = Inches(0.6); top = Inches(3.0)
cards = [
    ("Web app — Next.js",
     "Patients · scheduling · clinical · dental chart · treatment plans · prescriptions · "
     "invoicing · insurance · campaigns."),
    ("Mobile — React Native",
     "Clinician on the go: patient lookup, treatments, prescriptions, billing, WhatsApp inbox."),
    ("AI layer — Claude + GPT",
     "Real LLM features in production with prompt-cache cost engineering — not template substitution."),
]
for i, (lbl, body) in enumerate(cards):
    paragraph_card(s, left + i * (col_w + gap), top, col_w, col_h, lbl, body)

# Architecture ASCII strip → simple labeled boxes
add_textbox(s, Inches(0.6), Inches(5.3), Inches(12), Inches(0.3),
            "ARCHITECTURE", size=10, bold=True, color=MUTED)
arch_y = Inches(5.65)
box_w, box_h = Inches(2.2), Inches(0.7)
box_gap = Inches(0.2)
boxes = ["Next.js Web", "Expo Mobile", "NestJS API\n(multi-tenant)", "PostgreSQL\n+ Prisma", "AI Layer\n(Claude · GPT)"]
bx = Inches(0.6)
for label in boxes:
    add_card(s, bx, arch_y, box_w, box_h, fill=RGBColor(0x14, 0x22, 0x3D))
    add_textbox(s, bx, arch_y, box_w, box_h, label,
                size=11, bold=True, color=INK, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
    bx += box_w + box_gap

add_footer(s, 3)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Product surface
# ──────────────────────────────────────────────────────────────────────────────
s = add_slide()
section_label(s, "Product Surface")
slide_title(s, "Full clinic operations in one cloud-native suite.")

cols = [
    ("Clinical", ["Dental chart", "SOAP visits", "Treatment plans", "Prescriptions", "X-ray + photos"]),
    ("Operations", ["Smart scheduling", "Recall automation", "Inventory + stock", "Expenses", "Reports + exports"]),
    ("Revenue", ["Invoices + GST/tax", "Payments (Razorpay; Stripe Canada)", "Insurance claims", "Pre-auth + reimbursement", "Revenue insights"]),
    ("Engagement", ["WhatsApp Business", "SMS + Email", "Campaigns + A/B", "Referrals", "Feedback / NPS"]),
]
col_w = Inches(3.0); col_h = Inches(3.4); gap = Inches(0.15)
left = Inches(0.6); top = Inches(2.5)
for i, (lbl, bullets) in enumerate(cols):
    bullet_card(s, left + i * (col_w + gap), top, col_w, col_h, lbl, bullets)

add_textbox(s, Inches(0.6), Inches(6.15), Inches(12), Inches(0.6),
            "Phases 1–6 complete  ·  Phase 7 (Insurance/EHS) backend live, frontend in progress  ·  365+ backend tests  ·  Vitest + Playwright on frontend.",
            size=12, color=MUTED)
add_footer(s, 4)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — AI Relevance
# ──────────────────────────────────────────────────────────────────────────────
s = add_slide()
section_label(s, "AI Relevance — real AI, not basic automation")
slide_title(s, "Seven AI features in production today.")

ai_cards = [
    ("1 · AI Clinical Notes", "Generates structured SOAP notes; populates the same clinical model as manual entry — no parallel storage."),
    ("2 · AI Treatment Plan", "Multi-phase plans with urgency, alternatives, cost — full TreatmentPlan + Item rows."),
    ("3 · AI Prescription + Safety", "Allergy / age / pregnancy aware · drug-interaction flagged · inventory-aware (prescribes from clinic stock first)."),
    ("4 · AI Dental Chart Vision", "Claude Vision · GPT-4V — analyzes uploaded charts and X-rays (assistive, not diagnostic)."),
    ("5 · AI Revenue Insights", "Surfaces anomalies and growth levers from billing data."),
    ("6 · AI Patient Insights", "Nightly batch scores no-show, recall, churn, conversion. LLM reserved for click-to-explain. Zero hot-path AI cost."),
]
col_w = Inches(6.0); col_h = Inches(1.35); gap = Inches(0.2)
left1 = Inches(0.6); left2 = left1 + col_w + gap
top0 = Inches(2.4)
for i, (lbl, body) in enumerate(ai_cards):
    row = i // 2
    col = i % 2
    x = left1 if col == 0 else left2
    y = top0 + row * (col_h + gap)
    paragraph_card(s, x, y, col_w, col_h, lbl, body, body_size=12)

add_textbox(s, Inches(0.6), Inches(6.85), Inches(12), Inches(0.4),
            "+ AI Appointment Summary  ·  AI Campaign Content    ·    Voice-to-Consultation (Whisper) in roadmap.",
            size=11, color=MUTED)
add_footer(s, 5)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Productivity Impact
# ──────────────────────────────────────────────────────────────────────────────
s = add_slide()
section_label(s, "Productivity Impact")
slide_title(s, "Quantified targets, mapped to features.")

tiles = [
    ("Admin time", "−40%", "Recall automation, multi-channel reminders, AI-drafted messages."),
    ("No-shows", "−30%", "WhatsApp + SMS + risk scoring + waitlist auto-fill."),
    ("Documentation time", "−50%", "AI clinical notes + voice (roadmap)."),
    ("Operational efficiency", "+25%", "Unified UI + automation cron + dashboards."),
]
col_w = Inches(3.0); col_h = Inches(2.8); gap = Inches(0.15)
left = Inches(0.6); top = Inches(2.6)
for i, (lbl, val, sub) in enumerate(tiles):
    kpi_tile(s, left + i * (col_w + gap), top, col_w, col_h, lbl, val, sub)

add_textbox(s, Inches(0.6), Inches(5.85), Inches(12), Inches(0.5),
            "Targets to be validated via 3-clinic Canadian pilot in Year 1 (Section 9 of grant proposal).",
            size=12, color=MUTED)
add_footer(s, 6)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — Market
# ──────────────────────────────────────────────────────────────────────────────
s = add_slide()
section_label(s, "Market")
slide_title(s, "Canadian dental software · $[TBD]M TAM and growing.")

tiles = [
    ("TAM", "$[TBD]M", "Global dental practice mgmt software"),
    ("SAM", "$[TBD]M", "Canada + adjacent English markets"),
    ("SOM (Y3)", "$[TBD]M", "Target 1,500 Canadian clinics"),
]
col_w = Inches(2.5); col_h = Inches(2.2); gap = Inches(0.15)
left = Inches(0.6); top = Inches(2.4)
for i, (lbl, val, sub) in enumerate(tiles):
    kpi_tile(s, left + i * (col_w + gap), top, col_w, col_h, lbl, val, sub)

# Right column — playbook card
play_left = Inches(8.6); play_w = Inches(4.2); play_h = Inches(4.5)
paragraph_card(s, play_left, top, play_w, play_h, "Geographic playbook",
               "Year 1:  Ontario + BC pilots → Canada launch.\n\n"
               "Year 2:  Canada scale + select US states (insurance strategies for "
               "Delta Dental & MetLife already seeded in code).\n\n"
               "Year 3:  US expansion + UK / AU exploration. India pilot continues "
               "as cost-efficient revenue base.",
               body_size=13)

# Bottom market bullets
bullet_card(s, Inches(0.6), Inches(4.85), Inches(7.8), Inches(2.0),
            "Why Canada now",
            ["~16,000 dental practices in Canada",
             "Independent + DSO mix — both underserved by AI",
             "Incumbents priced at CAD $200–500+/op/month",
             "Cloud-native AI-first entrant has clear pricing + product wedge"],
            body_size=13)
add_footer(s, 7)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — Competitive landscape (table)
# ──────────────────────────────────────────────────────────────────────────────
s = add_slide()
section_label(s, "Competitive Landscape")
slide_title(s, "Bolt-on AI vs. AI-native cloud platform.")

# Build a real table
rows = [
    ["Capability", "Dentrix / ABEL / ClearDent", "Dentally / Curve", "Smart Dental Desk"],
    ["Architecture", "On-prem / hybrid", "Cloud", "Cloud-native, multi-tenant"],
    ["AI", "Bolted-on / partner", "Limited", "7 live features, prompt-cached"],
    ["Mobile (clinician)", "Limited", "Web-mobile", "Native React Native app"],
    ["WhatsApp / multi-channel", "No", "Partial", "Native + Meta-approved"],
    ["Insurance/EHS abstraction", "Region-locked", "Region-locked", "Country-strategy pattern"],
    ["Entry price", "$200–500+/op/mo", "$150–250/op/mo", "$49–149/clinic/mo"],
]
n_rows = len(rows); n_cols = len(rows[0])
tbl_left = Inches(0.6); tbl_top = Inches(2.4)
tbl_w = Inches(12.1); tbl_h = Inches(4.3)
table_shape = s.shapes.add_table(n_rows, n_cols, tbl_left, tbl_top, tbl_w, tbl_h)
table = table_shape.table
# Column widths
table.columns[0].width = Inches(2.6)
table.columns[1].width = Inches(3.2)
table.columns[2].width = Inches(3.1)
table.columns[3].width = Inches(3.2)
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = table.cell(ri, ci)
        cell.text = ""
        tf = cell.text_frame
        tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.08); tf.margin_bottom = Inches(0.08)
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = val
        r.font.name = "Calibri"
        if ri == 0:
            r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = MUTED
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x14, 0x22, 0x3D)
        else:
            r.font.size = Pt(12)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PANEL if ri % 2 else RGBColor(0x10, 0x1C, 0x33)
            if ci == 3:
                r.font.bold = True; r.font.color.rgb = ACCENT2
            elif ci == 0:
                r.font.color.rgb = INK; r.font.bold = True
            else:
                r.font.color.rgb = MUTED

add_footer(s, 8)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — Traction & Readiness
# ──────────────────────────────────────────────────────────────────────────────
s = add_slide()
section_label(s, "Traction & Readiness")
slide_title(s, "Production platform with paying clinics today.")

tiles = [
    ("Clinics live", "[X]", "Pilot market (India) — paying"),
    ("MRR", "₹[X]K", "Validates pricing + retention"),
    ("Code maturity", "365+", "Backend tests · Vitest + Playwright"),
    ("Mobile app", "v1", "Expo build, app-store ready"),
]
col_w = Inches(3.0); col_h = Inches(2.0); gap = Inches(0.15)
left = Inches(0.6); top = Inches(2.4)
for i, (lbl, val, sub) in enumerate(tiles):
    kpi_tile(s, left + i * (col_w + gap), top, col_w, col_h, lbl, val, sub)

# Two columns: done + next
bullet_card(s, Inches(0.6), Inches(4.65), Inches(6.2), Inches(2.2),
            "What's done",
            ["Phases 1–6: backend, web, hardening, comms, prod, AI, mobile",
             "Phase 7 (Insurance/EHS) backend complete — 11 Prisma models",
             "WhatsApp Meta-approved · embedded signup in flight"],
            body_size=12)

bullet_card(s, Inches(6.95), Inches(4.65), Inches(5.75), Inches(2.2),
            "What's next (with funding)",
            ["PIPEDA/PHIPA adaptation · Canadian cloud region",
             "3-clinic Canadian pilot · publish KPI study",
             "Public Canadian launch + association partnerships",
             "SOC 2 Type I certification"],
            label_color=ACCENT2, body_size=12)
add_footer(s, 9)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — Business Model
# ──────────────────────────────────────────────────────────────────────────────
s = add_slide()
section_label(s, "Business Model")
slide_title(s, "Validated 3-tier SaaS — repriced for Canada.")

rows = [
    ["Plan", "India (validated)", "Canada (proposed)", "Who it's for"],
    ["Free", "₹0", "CAD $0", "Try-before-buy · 1 doctor"],
    ["Standard ⭐", "₹699/mo", "CAD $49–69/mo", "Solo/duo practice · most popular"],
    ["Growth", "₹1,299/mo", "CAD $99–149/mo", "Multi-op · automation · own WABA"],
    ["Enterprise", "Custom", "Custom", "DSO / multi-branch chains"],
]
n_rows = len(rows); n_cols = len(rows[0])
tbl_shape = s.shapes.add_table(n_rows, n_cols, Inches(0.6), Inches(2.4),
                               Inches(12.1), Inches(2.8))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(2.4)
tbl.columns[1].width = Inches(2.8)
tbl.columns[2].width = Inches(3.0)
tbl.columns[3].width = Inches(3.9)
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.text = ""
        tf = cell.text_frame
        tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.06); tf.margin_bottom = Inches(0.06)
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = val
        r.font.name = "Calibri"
        if ri == 0:
            r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = MUTED
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x14, 0x22, 0x3D)
        else:
            r.font.size = Pt(13)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PANEL if ri % 2 else RGBColor(0x10, 0x1C, 0x33)
            if ci == 0:
                r.font.bold = True; r.font.color.rgb = INK
            elif ci == 2:
                r.font.bold = True; r.font.color.rgb = ACCENT2
            else:
                r.font.color.rgb = INK if ci == 1 else MUTED

# Bottom row — 3 KPI tiles
col_w = Inches(4.0); col_h = Inches(1.6); gap = Inches(0.15)
left = Inches(0.6); top = Inches(5.4)
econ = [
    ("Gross margin", "80%+", "Cloud SaaS, capped AI quota"),
    ("Overage", "PAYG", "WhatsApp / SMS / AI usage"),
    ("Expansion", "Marketplace", "Labs · imaging · insurance fees (Y2+)"),
]
for i, (lbl, val, sub) in enumerate(econ):
    kpi_tile(s, left + i * (col_w + gap), top, col_w, col_h, lbl, val, sub)

add_footer(s, 10)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — Collaboration / Consortium
# ──────────────────────────────────────────────────────────────────────────────
s = add_slide()
section_label(s, "Collaboration & Consortium")
slide_title(s, "Building a Canadian healthcare-AI alliance.")

cards = [
    ("Academic partner",
     "Canadian dental school for clinical validation, research publication, and resident training. [Schulich / U of T / UBC — outreach planned]"),
    ("Clinical partner",
     "Multi-location DSO or independent dentist association for pilot deployment and real-world data. [3 candidates being identified]"),
    ("Cloud + infra partner",
     "AWS Canada Central / Azure Canada Central, MLOps support, Canadian data residency. [CSP engagement in progress]"),
    ("Compliance partner",
     "PIPEDA / PHIPA specialist for initial audit and ongoing compliance. [Quote in progress]"),
]
col_w = Inches(6.0); col_h = Inches(1.9); gap = Inches(0.2)
left1 = Inches(0.6); left2 = left1 + col_w + gap
top0 = Inches(2.4)
for i, (lbl, body) in enumerate(cards):
    row = i // 2; col = i % 2
    x = left1 if col == 0 else left2
    y = top0 + row * (col_h + gap)
    paragraph_card(s, x, y, col_w, col_h, lbl, body, body_size=12)

add_textbox(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.5),
            "LOIs and MOUs to be attached as Annex A of the grant submission.",
            size=12, color=MUTED)
add_footer(s, 11)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — Team
# ──────────────────────────────────────────────────────────────────────────────
s = add_slide()
section_label(s, "Team")
slide_title(s, "Builders who ship.")

team = [
    ("Founder / CEO", "[Name]", "[Background — clinical / commercial leadership]"),
    ("Co-founder / Tech lead", "Prasanth V",
     "Full-stack engineer. Built the entire platform across web, backend, mobile, AI integration, and prompt-cache cost engineering."),
    ("Clinical advisor", "[Name, BDS/MDS]",
     "Practicing dentist · clinical workflow validation · advisory board."),
]
col_w = Inches(4.0); col_h = Inches(2.6); gap = Inches(0.15)
left = Inches(0.6); top = Inches(2.4)
for i, (role, name, bio) in enumerate(team):
    x = left + i * (col_w + gap)
    add_card(s, x, top, col_w, col_h)
    add_textbox(s, x + Inches(0.25), top + Inches(0.2),
                col_w - Inches(0.5), Inches(0.3),
                role.upper(), size=10, bold=True, color=ACCENT)
    add_textbox(s, x + Inches(0.25), top + Inches(0.55),
                col_w - Inches(0.5), Inches(0.6),
                name, size=20, bold=True, color=INK)
    add_textbox(s, x + Inches(0.25), top + Inches(1.2),
                col_w - Inches(0.5), col_h - Inches(1.35),
                bio, size=12, color=MUTED, line_spacing=1.3)

# Bottom row — 3 helper cards
col_w = Inches(4.0); col_h = Inches(1.6); gap = Inches(0.15)
top = Inches(5.25)
helpers = [
    ("Post-funding hires", "Canadian market lead · 2nd backend engineer · Customer success"),
    ("Advisors to add", "Practicing Canadian dentist · Healthcare regulatory veteran · GTM lead"),
    ("Why us", "We've already built and shipped what most pre-seed teams promise. Funding accelerates a working product into a new geography."),
]
for i, (lbl, body) in enumerate(helpers):
    paragraph_card(s, Inches(0.6) + i * (col_w + gap), top, col_w, col_h, lbl, body, body_size=12)
add_footer(s, 12)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 13 — Ask & Use of Funds
# ──────────────────────────────────────────────────────────────────────────────
s = add_slide()
section_label(s, "The Ask")
add_rich(s, Inches(0.6), Inches(0.95), Inches(12), Inches(1.5), [
    ("CAD $[TBD]", {"size": 40, "bold": True, "color": ACCENT2}),
    ("  over  ", {"size": 32, "color": INK}),
    ("[TBD] months", {"size": 40, "bold": True, "color": ACCENT2}),
    (".", {"size": 32, "color": INK}),
], line_spacing=1.1)

# Use of funds table
rows = [
    ["Use of funds", "%"],
    ["Engineering (Canadian features, compliance, hardening)", "35%"],
    ["Pilot operations + clinical validation", "20%"],
    ["Go-to-market (marketing, content, sales)", "20%"],
    ["Compliance + security (SOC 2, PIPEDA, pen test)", "10%"],
    ["Infrastructure (Canadian cloud)", "5%"],
    ["Working capital + contingency", "10%"],
]
n_rows = len(rows); n_cols = 2
tbl_shape = s.shapes.add_table(n_rows, n_cols, Inches(0.6), Inches(3.0),
                               Inches(7.4), Inches(3.7))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(6.0); tbl.columns[1].width = Inches(1.4)
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.text = ""
        tf = cell.text_frame
        tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.06); tf.margin_bottom = Inches(0.06)
        p = tf.paragraphs[0]
        if ci == 1:
            p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = val
        r.font.name = "Calibri"
        if ri == 0:
            r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = MUTED
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x14, 0x22, 0x3D)
        else:
            r.font.size = Pt(13); r.font.color.rgb = INK
            if ci == 1:
                r.font.bold = True; r.font.color.rgb = ACCENT2
            cell.fill.solid()
            cell.fill.fore_color.rgb = PANEL if ri % 2 else RGBColor(0x10, 0x1C, 0x33)

# Milestones card on the right
bullet_card(s, Inches(8.4), Inches(3.0), Inches(4.3), Inches(3.7),
            "12-month milestones",
            ["M1: Canadian cloud live · PIPEDA gap analysis",
             "M3: First pilot clinic live in Ontario",
             "M6: Validated KPIs from 3 pilots",
             "M7: Public Canadian launch",
             "M11: SOC 2 Type I complete",
             "M12: 50 paying Canadian clinics"],
            label_color=ACCENT2, body_size=12)
add_footer(s, 13)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 14 — Closing
# ──────────────────────────────────────────────────────────────────────────────
s = add_slide()
add_pill(s, Inches(0.6), Inches(1.4), "LET'S BUILD IT")
add_rich(s, Inches(0.6), Inches(2.0), Inches(12), Inches(3.5), [
    ("A scalable Canadian", {"size": 50, "bold": True, "color": INK}),
    ("\n", {}),
    ("healthcare-AI platform —", {"size": 50, "bold": True, "color": INK}),
    ("\n", {}),
    ("shipped, not promised.", {"size": 50, "bold": True, "color": ACCENT2}),
], line_spacing=1.08)

add_textbox(s, Inches(0.6), Inches(5.2), Inches(11), Inches(1.0),
            "We have the product, the architecture, and the pilot data path. With this round we bring it to Canadian dental practices "
            "and a measurable, publishable clinical outcome.",
            size=16, color=MUTED, line_spacing=1.4)

# Contact chips
x = Inches(0.6)
for c in ["pvprasanth474@gmail.com", "smartdentaldesk.com", "Smart Dental Desk · FB 61572113176366"]:
    w = add_chip(s, x, Inches(6.4), c)
    x += w + Inches(0.15)

add_textbox(s, Inches(0.6), Inches(6.85), Inches(12), Inches(0.4),
            "→  Schedule a 30-minute product walkthrough.",
            size=14, bold=True, color=ACCENT2)
add_footer(s, 14)


# ── Save ────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(__file__), "SmartDentalDesk-PitchDeck.pptx")
prs.save(out)
print(f"Saved: {out}")
